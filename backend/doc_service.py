from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt
import io

def create_word_doc(title: str, structure: list):
    doc = Document()
    doc.add_heading(title, 0)

    for section in structure:
        # Add Header
        doc.add_heading(section['title'], level=1)
        # Add Content
        # (Simple cleaning to remove markdown symbols if necessary)
        clean_content = section['content'].replace('**', '').replace('#', '')
        doc.add_paragraph(clean_content)
        
    # Save to memory buffer
    buffer = io.BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    return buffer

def create_ppt_doc(title: str, structure: list):
    prs = Presentation()
    
    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = title
    
    for section in structure:
        # Bullet Layout
        bullet_slide_layout = prs.slide_layouts[1] 
        slide = prs.slides.add_slide(bullet_slide_layout)
        
        # Set Title
        slide.shapes.title.text = section['title']
        
        # Set Content
        tf = slide.shapes.placeholders[1].text_frame
        tf.text = section['content'].replace('**', '') # Simple cleanup

    # Save to memory buffer
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer